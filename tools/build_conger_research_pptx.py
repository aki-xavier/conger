#!/usr/bin/env python3
"""可重复生成 conger 研究汇报 PPTX。

主题: 从 Riesz-GMM 显式建模 到 conger 当前实例级 MixtureSPN 逆渲染。

运行 (项目 .venv 已临时安装 python-pptx, 不改 pyproject/uv.lock):
    .venv/bin/python tools/build_conger_research_pptx.py

产出:
    docs/conger_riesz_gmm_to_mixture_spn.pptx

本脚本同时内置结构验收函数 validate()——重新打开产物, 校验 16:9、
页数、shape 边界与文本溢出启发式。全部证据为 git/docs 核对后的常量。
"""

from __future__ import annotations

import os
import unicodedata
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ───────────────────────────────────────────────────────────────────────────
# 常量
# ───────────────────────────────────────────────────────────────────────────

REPO = Path(__file__).resolve().parent.parent
OUT_PPTX = REPO / "docs" / "conger_riesz_gmm_to_mixture_spn.pptx"
OUTLINE_MD = REPO / "docs" / "conger_riesz_gmm_to_mixture_spn_outline.md"

ART_SCATTER = REPO / "artifacts" / "inverse_scatter.png"
ART_RECON = REPO / "artifacts" / "inverse_recon.png"

SLIDE_W = Emu(12_192_000)  # 13.333 in
SLIDE_H = Emu(6_858_000)   # 7.5 in
SLIDE_W_IN = 13.333333333333334
SLIDE_H_IN = 7.5

FONT = "Noto Sans CJK SC"

# 安静研究风配色
DARK = RGBColor(0x18, 0x31, 0x53)   # 深蓝  #183153
TEAL = RGBColor(0x0F, 0x76, 0x6E)   # 青绿  #0F766E
RED = RGBColor(0xD6, 0x45, 0x50)    # 失败/否决 红
GOLD = RGBColor(0xC7, 0x8A, 0x22)   # 证据/保留 金
LIGHT = RGBColor(0xF4, 0xF5, 0xF7)  # 浅灰底
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x69, 0x71, 0x7E)   # 弱化文字
MID = RGBColor(0xDC, 0xE0, 0xE6)    # 分隔线/浅边框
BLUE_SOFT = RGBColor(0xDD, 0xE6, 0xF0)
TEAL_SOFT = RGBColor(0xD6, 0xEB, 0xE9)
RED_SOFT = RGBColor(0xF6, 0xDD, 0xDF)
GOLD_SOFT = RGBColor(0xF3, 0xE7, 0xD2)

TOTAL = 15


# ───────────────────────────────────────────────────────────────────────────
# 低层原语
# ───────────────────────────────────────────────────────────────────────────

def _ea_width(ch: str) -> str:
    return unicodedata.east_asian_width(ch)


def _char_w_in(ch: str, size_in: float) -> float:
    """单字符近似宽度 (英寸)。CJK 全宽 ≈ 1×字号, 其余 ≈ 0.52×。"""
    if ch == " ":
        return 0.28 * size_in
    if _ea_width(ch) in ("F", "W"):
        return size_in
    return 0.52 * size_in


def run(text: str, size: float, bold: bool = False, color: RGBColor = DARK,
        italic: bool = False) -> tuple:
    return (text, size, bold, color, italic)


def _no_shadow(shape) -> None:
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def _set_run_font(r, size: float, bold: bool, color: RGBColor, italic: bool) -> None:
    """同时设 latin / ea / cs 三路 typeface, 确保中文走 Noto Sans CJK SC。"""
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = FONT  # 写 <a:latin>
    rPr = r._r.get_or_add_rPr()
    # 追加 <a:ea> 与 <a:cs>
    from pptx.oxml.ns import qn
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)


def tx(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
       wrap=True, space_after=2.0) -> None:
    """paras: list[ list[run-tuple] ]; 每个内层 list 是一段, 元素为 run() 元组。"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(space_after)
        p.line_spacing = 1.0
        for (text, size, bold, color, italic) in para:
            r = p.add_run()
            r.text = text
            _set_run_font(r, size, bold, color, italic)
    return box


def shape(slide, x, y, w, h, fill, line=None, kind=MSO_SHAPE.RECTANGLE,
          line_w=0.75, radius=None) -> object:
    sp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    _no_shadow(sp)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    if radius is not None and kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def hline(slide, x, y, w, color=MID, weight=0.75) -> None:
    shape(slide, x, y, w, 0.012, color, line=color, line_w=weight)


def footer(slide, num: int, evidence: str) -> None:
    hline(slide, 0.5, 7.06, 12.333, MID, 1.0)
    tx(slide, 0.5, 7.14, 10.0, 0.26, [[run(evidence, 9, False, GRAY)]])
    tx(slide, 11.3, 7.14, 1.53, 0.26, [[run(f"{num} / {TOTAL}", 9, False, GRAY)]],
       align=PP_ALIGN.RIGHT)


def title(slide, text: str, kicker: str | None = None, accent: RGBColor = TEAL,
          evidence: str | None = None) -> None:
    shape(slide, 0.5, 0.5, 0.14, 0.82, accent)
    if kicker:
        tx(slide, 0.82, 0.38, 11.8, 0.3, [[run(kicker, 11, True, accent)]])
        tx(slide, 0.82, 0.66, 11.8, 0.72, [[run(text, 28, True, DARK)]],
           anchor=MSO_ANCHOR.MIDDLE)
    else:
        tx(slide, 0.82, 0.44, 11.8, 0.9, [[run(text, 28, True, DARK)]],
           anchor=MSO_ANCHOR.MIDDLE)
    if evidence:
        tx(slide, 0.5, 1.38, 12.3, 0.26, [[run(evidence, 9.5, False, GRAY)]])


def chip(slide, x, y, w, h, text, fill, fg=WHITE, size=12, bold=True,
         radius=0.5) -> None:
    shape(slide, x, y, w, h, fill, radius=radius, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
    tx(slide, x, y, w, h, [[run(text, size, bold, fg)]],
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)


def node(slide, x, y, w, h, text, fill=BLUE_SOFT, fg=DARK, size=12.5,
         bold=True, sub: str | None = None, line=TEAL) -> None:
    shape(slide, x, y, w, h, fill, line=line, line_w=1.0,
          kind=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
    if sub:
        tx(slide, x + 0.08, y + 0.05, w - 0.16, h - 0.1,
           [[run(text, size, bold, fg)], [run(sub, size - 2.5, False, GRAY)]],
           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=1)
    else:
        tx(slide, x + 0.06, y, w - 0.12, h, [[run(text, size, bold, fg)]],
           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)


def arrow_r(slide, x, y, w, h=0.34, color=TEAL) -> None:
    shape(slide, x, y, w, h, color, kind=MSO_SHAPE.RIGHT_ARROW)


def arrow_d(slide, x, y, w, h, color=TEAL) -> None:
    shape(slide, x, y, w, h, color, kind=MSO_SHAPE.DOWN_ARROW)


def _fit_image(slide, path, x, y, max_w, max_h) -> tuple[float, float]:
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    ar = ih / iw
    w = max_w
    h = w * ar
    if h > max_h:
        h = max_h
        w = h / ar
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    return w, h


# ───────────────────────────────────────────────────────────────────────────
# 各页
# ───────────────────────────────────────────────────────────────────────────

def s1_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    shape(s, 0, 0, SLIDE_W_IN, SLIDE_H_IN, WHITE)
    shape(s, 0, 0, SLIDE_W_IN, 0.18, TEAL)
    shape(s, 0, SLIDE_H_IN - 0.18, SLIDE_W_IN, 0.18, DARK)

    tx(s, 0.9, 2.0, 11.5, 0.4, [[run("CONGER · 研究汇报 · 2026-08", 13, True, TEAL)]])
    tx(s, 0.9, 2.45, 11.6, 1.9, [
        [run("从 Riesz-GMM 显式建模", 40, True, DARK)],
        [run("到实例级 MixtureSPN 逆渲染", 40, True, DARK)],
    ], space_after=6)
    hline(s, 0.9, 4.35, 2.2, TEAL, 2.0)
    tx(s, 0.9, 4.6, 11.0, 0.9, [
        [run("一条由实测否决驱动的路线切换:", 17, False, GRAY)],
        [run("模块显式 → 生成过程显式 + 后验学习", 17, True, DARK)],
    ], space_after=3)
    tx(s, 0.9, 5.6, 11.0, 0.5, [
        [run("左右立体图  →  Riesz 全分辨率特征  →  完整 cga.Scene(含光照)", 14, False, GRAY)],
    ])
    footer(s, 1, "8ed8356 (07-24) · aaa43e0 (07-27) · 97c33d6 (08-11) · 523b97a / 54cd2d1 (08-13)")
    return s


def s2_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "一句话总览: 从「模块显式」到「生成过程显式 + 后验学习」",
          kicker="研究定位", evidence="docs/architecture.md §0")
    # 左: 旧
    shape(s, 0.6, 1.75, 5.9, 2.9, LIGHT, radius=0.06, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
    tx(s, 0.85, 1.95, 5.4, 0.4, [[run("旧: 模块显式", 18, True, DARK)]])
    tx(s, 0.85, 2.45, 5.45, 2.0, [
        [run("逐模块闭式代数", 15, False, DARK)],
        [run("手工先验 / 门控 / 参数 β·T", 15, False, DARK)],
        [run("无训练数据 · 无梯度优化", 15, False, DARK)],
        [run("软隶属度 Q 不流入深度", 15, True, RED)],
    ], space_after=5)
    # 中: 箭头
    arrow_r(s, 6.6, 3.0, 0.8, 0.4, TEAL)
    # 右: 新
    shape(s, 7.5, 1.75, 5.4, 2.9, TEAL_SOFT, radius=0.06, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
    tx(s, 7.75, 1.95, 4.9, 0.4, [[run("新: 生成过程显式 + 后验学习", 18, True, DARK)]])
    tx(s, 7.75, 2.45, 4.95, 2.0, [
        [run("单一生成过程: Codebook → render", 15, False, DARK)],
        [run("统计规律交给实例级后验", 15, False, DARK)],
        [run("硬约束进支持集 · 歧义进候选后验", 15, False, DARK)],
        [run("实测驱动 · 诚实否决", 15, True, TEAL)],
    ], space_after=5)
    # 底部: 三阶段时间轴
    tx(s, 0.6, 5.0, 6.0, 0.32, [[run("路线时间轴", 12, True, GRAY)]])
    yl = 5.4
    hline(s, 0.6, yl + 0.06, 12.1, MID, 2.0)
    segs = [
        (0.6, "显式建模期", "07-24 → 08-05", "8ed8356 · 3a344b4 · cbc9fa8 · 43416c8", DARK),
        (4.7, "SPN 分支验证", "08-11", "285763a · 97c33d6", GOLD),
        (8.2, "实例级主线", "08-13", "523b97a · 54cd2d1", TEAL),
    ]
    for (x, lab, d, ev, c) in segs:
        shape(s, x + 0.14, yl, 0.16, 0.16, c, kind=MSO_SHAPE.OVAL)
        tx(s, x - 0.1, yl + 0.26, 3.4, 0.3, [[run(lab, 13, True, c)]])
        tx(s, x - 0.1, yl + 0.54, 3.4, 0.26, [[run(d, 10.5, False, GRAY)]])
        tx(s, x - 0.1, yl + 0.8, 3.4, 0.26, [[run(ev, 9, False, GRAY)]])
    footer(s, 2, "docs/architecture.md §0 · git log")
    return s


def s3_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "当前要解的逆渲染问题", kicker="问题定义",
          evidence="docs/prior.md · docs/architecture.md §1")
    # 输入 → 输出
    node(s, 0.7, 1.9, 3.4, 1.5, "输入", sub="左右两张 144×144 立体图", fill=BLUE_SOFT)
    arrow_r(s, 4.15, 2.45, 0.7, 0.4, TEAL)
    node(s, 4.9, 1.9, 3.6, 1.5, "?", sub="不适定: 无穷多场景同图", fill=RED_SOFT, line=RED)
    arrow_r(s, 8.55, 2.45, 0.7, 0.4, TEAL)
    node(s, 9.3, 1.9, 3.4, 1.5, "输出", sub="完整 cga.Scene(含光照)", fill=TEAL_SOFT)
    # 输出自由度
    tx(s, 0.7, 3.75, 6.0, 0.32, [[run("输出自由度(单图元)", 12, True, GRAY)]])
    dofs = ["kind 图元", "u,v 位置", "s 尺寸", "z 深度", "hue 色相", "lcol 光色", "ldir 光向"]
    x = 0.7
    for d in dofs:
        chip(s, x, 4.12, 1.55, 0.42, d, DARK, size=11)
        x += 1.68
    # 先验形式
    tx(s, 0.7, 4.95, 6.0, 0.32, [[run("不适定 → 需要强先验(Helmholtz 无意识推理)", 12, True, GRAY)]])
    shape(s, 0.7, 5.32, 12.0, 1.5, LIGHT, radius=0.05, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
    tx(s, 0.95, 5.5, 11.5, 1.2, [
        [run("本项目的先验四层落地:  ", 13.5, True, DARK),
         run("场景族硬先验 · 经验统计先验(实例记忆/均匀权重/tied方差) · 算法结构先验 · 推理似然先验", 13.5, False, DARK)],
        [run("设计原则: 硬物理约束进支持集, 统计规律进概率/核带宽, 不适定歧义进候选后验, 不做过早 argmax。", 13.5, False, DARK)],
    ], space_after=4)
    footer(s, 3, "docs/prior.md · docs/architecture.md §1")
    return s


def s4_pipeline(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "当前主链路架构", kicker="架构总图",
          evidence="docs/architecture.md §1 · 模块结构 §5")
    row1 = [
        ("左右立体图", "144×144 ×2"),
        ("Riesz 特征", "11 通道 · V≈228K"),
        ("PCA 白化", "Gram eigh · D≤N−1"),
        ("MixtureSPN", "逐 kind · K=N · 无 EM"),
    ]
    row2 = [
        ("条件期望/后验", "连续=核回归 · 离散=后验"),
        ("渲染残差精炼", "top-k kind × 54 外观"),
        ("StructuredHypothesis", "MAP + 候选 + 联合后验"),
        ("cga.Scene", "完整重建(含光照)"),
    ]
    x = 0.55
    w = 2.92
    for i, (t, sub) in enumerate(row1):
        node(s, x, 1.85, w, 1.25, t, sub=sub, fill=BLUE_SOFT)
        if i < 3:
            arrow_r(s, x + w + 0.02, 2.3, 0.24, 0.3, TEAL)
        x += w + 0.28
    # 下行右→左视觉上接续: 用一条下行+反向箭头示意
    arrow_d(s, 6.35, 3.12, 0.3, 0.4, TEAL)
    # 第二行 (从左到右, 与第一行末接续)
    x = 0.55
    for i, (t, sub) in enumerate(row2):
        node(s, x, 3.75, w, 1.25, t, sub=sub, fill=TEAL_SOFT if i >= 1 else BLUE_SOFT)
        if i < 3:
            arrow_r(s, x + w + 0.02, 4.2, 0.24, 0.3, TEAL)
        x += w + 0.28
    # 关键说明
    tx(s, 0.7, 5.4, 12.0, 1.4, [
        [run("要点:  ", 14, True, DARK),
         run("单一生成过程(Codebook → cga 渲染)取代逐模块手工先验; 白化使对角高斯 ≡ 原空间全协方差; 实例级组件 = 每个训练样本一个分量, 无 EM/无质心压缩。", 14, False, DARK)],
        [run("渲染残差精炼把反照率×光照歧义交回正向模型 —— 这是外观辨识的主要来源。", 14, False, DARK)],
    ], space_after=5)
    footer(s, 4, "docs/architecture.md §1 · §5")
    return s


def s5_old_pipeline(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "旧 Riesz-GMM 显式管线: 显式建模了什么", kicker="旧路线",
          evidence="aaa43e0 abandonment report")
    mods = ["知觉组织", "场景分割 / BN", "单目深度", "双目视差", "深度融合", "Motor-EKF", "CGA 重建"]
    x = 0.55
    w = 1.62
    for i, m in enumerate(mods):
        node(s, x, 1.85, w, 0.95, m, fill=LIGHT, fg=DARK, size=11.5, line=MID)
        if i < len(mods) - 1:
            arrow_r(s, x + w + 0.01, 2.1, 0.16, 0.28, GRAY)
        x += w + 0.185
    tx(s, 0.7, 3.15, 12.0, 0.34, [[run("三件套: Riesz 多尺度谱特征 + GMM/VB-GMM 软聚类 + CGA 图元场景图", 15, True, DARK)]])
    tx(s, 0.7, 3.6, 12.0, 1.1, [
        [run("逐模块闭式、手工先验/门控、参数 β·T 标定;", 15, False, DARK),
         run("无训练数据、无梯度优化、可审计、可复用", 15, True, DARK),
         run("—— 四个硬约束下的自洽设计。", 15, False, DARK)],
        [run("旧 SPN 早期场景码: 288 = kind3 × 8×6 位置网格 × 2 尺寸 —— 连续物理量被离散化。", 14.5, False, DARK)],
    ], space_after=5)
    # 时间轴(旧路线关键提交)
    tx(s, 0.7, 4.85, 8.0, 0.32, [[run("旧路线关键提交", 12, True, GRAY)]])
    yl = 5.25
    hline(s, 0.7, yl + 0.05, 12.0, MID, 2.0)
    commits = [
        (0.7, "8ed8356", "07-24 初始", DARK),
        (2.9, "aaa43e0", "07-27 废弃报告(DGGG/DG-MRF)", RED),
        (5.2, "b52ac7a", "07-31 CGA implementation", DARK),
        (7.3, "3a344b4", "08-02 add riesz.py", DARK),
        (9.4, "cbc9fa8", "08-04 remove gmm.py", DARK),
        (11.4, "43416c8", "08-05 flow/params", DARK),
    ]
    for (cx, h, lab, c) in commits:
        shape(s, cx + 0.05, yl, 0.12, 0.12, c, kind=MSO_SHAPE.OVAL)
        tx(s, cx - 0.25, yl + 0.18, 1.9, 0.26, [[run(h, 10.5, True, c)]])
        tx(s, cx - 0.25, yl + 0.44, 1.9, 0.4, [[run(lab, 9, False, GRAY)]])
    footer(s, 5, "aaa43e0 · b52ac7a · 3a344b4 · cbc9fa8 · 43416c8")
    return s


def s6_kept(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "旧路线的价值: 哪些组件被保留", kicker="旧路线",
          evidence="aaa43e0 abandonment report")
    kept = [
        ("Riesz 前端", "多尺度谱特征"),
        ("2θ 圆统计", "方向用圆分布而非高斯"),
        ("复数色相 S·e^{iH}", "等亮度色度编码"),
        ("增益控制", "对比度/光照不变"),
        ("合成信号测试", "9 信号标准测试集"),
        ("CGA 场景表示", "图元场景图"),
        ("双目几何 / 前景权重", "立体锚点"),
    ]
    x = 0.7
    y = 1.85
    w = 3.9
    h = 1.0
    for i, (t, sub) in enumerate(kept):
        col = i % 3
        row = i // 3
        gx = 0.7 + col * 4.15
        gy = 1.85 + row * 1.3
        shape(s, gx, gy, 3.85, 1.0, GOLD_SOFT, line=GOLD, line_w=1.0,
              kind=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.16)
        tx(s, gx + 0.15, gy + 0.1, 3.55, 0.34, [[run(t, 14, True, DARK)]])
        tx(s, gx + 0.15, gy + 0.48, 3.55, 0.4, [[run(sub, 11, False, GRAY)]])
    tx(s, 0.7, 6.0, 12.0, 0.6, [
        [run("结论:  ", 15, True, GOLD),
         run("失败的只是「聚类 + 手工先验」的显式推理层; 信号前端与几何表示是正确资产, 被后续 SPN 路线直接继承。", 15, False, DARK)],
    ], space_after=3)
    footer(s, 6, "aaa43e0 · docs/architecture.md §2.1")
    return s


def s7_fail1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "失败证据 I: DGGG → DG-MRF 的标定困境", kicker="旧路线否决",
          evidence="aaa43e0 (DGGG/DG-MRF 数值)")
    # 原生条形图
    groups = [("quadrants", 85.7, 93.1), ("orientation", 54.1, 97.9), ("shaded", 78.8, 93.8)]
    base_x, base_y = 0.8, 1.9
    max_h = 3.3
    gw = 2.05
    barw = 0.62
    # 刻度与网格
    for gval in (50, 100):
        gy = base_y + max_h * (1 - gval / 100)
        hline(s, base_x - 0.15, gy, 6.6, MID, 0.5)
        tx(s, base_x - 0.42, gy - 0.1, 0.4, 0.2, [[run(str(gval), 9, False, GRAY)]], align=PP_ALIGN.RIGHT)
    for gi, (lab, dggg, dmrf) in enumerate(groups):
        cx = base_x + gi * gw
        tx(s, cx + 0.1, base_y + max_h + 0.08, gw, 0.26, [[run(lab, 11, True, DARK)]], align=PP_ALIGN.CENTER)
        # DGGG 柱
        hd = max_h * dggg / 100
        col = RED if gi == 1 else BLUE_SOFT
        shape(s, cx, base_y + max_h - hd, barw, hd, col, line=DARK if gi == 1 else None, line_w=1.0)
        tx(s, cx, base_y + max_h - hd - 0.24, barw, 0.22, [[run(f"{dggg}", 10, True, DARK)]], align=PP_ALIGN.CENTER)
        # DG-MRF 柱
        hm = max_h * dmrf / 100
        shape(s, cx + barw + 0.14, base_y + max_h - hm, barw, hm, TEAL)
        tx(s, cx + barw + 0.14, base_y + max_h - hm - 0.24, barw, 0.22, [[run(f"{dmrf}", 10, True, TEAL)]], align=PP_ALIGN.CENTER)
    # 图例
    shape(s, base_x, 5.85, 0.24, 0.24, RED)
    tx(s, base_x + 0.32, 5.85, 1.6, 0.24, [[run("DGGG", 11, False, DARK)]])
    shape(s, base_x + 2.0, 5.85, 0.24, 0.24, TEAL)
    tx(s, base_x + 2.32, 5.85, 1.6, 0.24, [[run("DG-MRF", 11, False, DARK)]])
    # 右侧: 标定困境
    rx = 7.6
    shape(s, rx, 1.9, 5.2, 4.2, LIGHT, radius=0.05, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
    tx(s, rx + 0.25, 2.05, 4.7, 0.4, [[run("MRF 修复精度, 但引入三参数标定", 14.5, True, RED)]])
    tx(s, rx + 0.25, 2.55, 4.75, 3.3, [
        [run("β", 14, True, DARK), run("  一票邻域 = 多少 nat —— 需逐场景网格搜索", 13, False, DARK)],
        [run("T", 14, True, DARK), run("  softmax 温度 —— 二分搜索 ~100–200 步 × 6 轮", 13, False, DARK)],
        [run("门控", 14, True, DARK), run("  Pb>0.5 两步后处理, 硬阈值 0.49/0.51 天差地别", 13, False, DARK)],
    ], space_after=5)
    tx(s, rx + 0.25, 5.1, 4.75, 0.9, [
        [run("12 路改进全部否决: ", 12.5, True, RED)],
        [run("HMRF-EM 61–66% · pyramid 87.2→48.5 · 显著性加权 53%", 12, False, DARK)],
    ], space_after=3)
    footer(s, 7, "aaa43e0 · DGGG: orientation 54.1 / shaded 78.8 · DG-MRF: orientation 97.9 / shaded 93.8")
    return s


def s8_fail2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "失败证据 II: 拓扑错误 + 单目/遮挡白卷", kicker="旧路线否决",
          evidence="aaa43e0 · fffc06a · 05ac7f0")
    cards = [
        ("信号流拓扑错误", "聚类占 640×480 热帧 37% / 2.5s, 但软隶属度 Q 不流入深度 —— 输出对深度精度贡献为零。", RED, RED_SOFT,
         "拓扑错误, 不是性能问题", "aaa43e0"),
        ("单目深度白卷", "iBims / NYU 室内: Spearman ≈ 0", RED, RED_SOFT,
         "单目线索域内失败", "fffc06a"),
        ("遮挡序白卷", "iBims 序正确率 0.444 ≈ 随机; 修正后 0.419; T 结假阳性 76%", RED, RED_SOFT,
         "序数约束不成立", "05ac7f0"),
    ]
    for i, (t, body, tag, bg, note, ev) in enumerate(cards):
        gy = 1.85 + i * 1.62
        shape(s, 0.7, gy, 12.0, 1.42, bg, radius=0.05, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
        shape(s, 0.7, gy, 0.14, 1.42, tag)
        tx(s, 1.05, gy + 0.14, 3.0, 0.4, [[run(t, 15, True, tag)]])
        tx(s, 1.05, gy + 0.6, 8.2, 0.7, [[run(body, 13.5, False, DARK)]], space_after=2)
        chip(s, 9.7, gy + 0.5, 2.7, 0.44, note, tag, size=10.5)
        tx(s, 9.7, gy + 1.0, 2.7, 0.26, [[run(ev, 9, False, GRAY)]], align=PP_ALIGN.CENTER)
    tx(s, 0.7, 6.75, 12.0, 0.4, [
        [run("共同根因: 手工先验/显式推理层产出与下游需求脱节 —— 转向「生成过程显式」的直接动因。", 13.5, True, DARK)],
    ])
    footer(s, 8, "aaa43e0 (37%/2.5s) · fffc06a (Spearman≈0) · 05ac7f0 (0.444→0.419 / T结 76%)")
    return s


def s9_pivot(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "转向的决定性实验: SPN 场景码 0.970", kicker="SPN 分支验证",
          evidence="97c33d6 (08-11)")
    # 大数字
    shape(s, 0.9, 2.0, 5.4, 3.3, DARK, radius=0.08, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
    tx(s, 1.1, 2.35, 5.0, 0.4, [[run("场景码整体准确率", 13, True, TEAL)]])
    tx(s, 1.1, 2.9, 5.0, 1.4, [[run("0.970", 72, True, WHITE)]])
    tx(s, 1.1, 4.35, 5.0, 0.7, [[run("N=2000 · 288 码枚举后验 argmax", 12.5, False, WHITE)]])
    # 各因子
    tx(s, 6.7, 2.0, 5.8, 0.34, [[run("各因子准确率", 14, True, DARK)]])
    facs = [("kind", "1.0"), ("gx", "0.99"), ("gy", "0.98"), ("size", "1.0")]
    x = 6.7
    for i, (k, v) in enumerate(facs):
        gx = 6.7 + i * 1.5
        shape(s, gx, 2.4, 1.32, 1.0, LIGHT, radius=0.12, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
        tx(s, gx, 2.52, 1.32, 0.3, [[run(k, 11, False, GRAY)]], align=PP_ALIGN.CENTER)
        tx(s, gx, 2.85, 1.32, 0.42, [[run(v, 20, True, TEAL)]], align=PP_ALIGN.CENTER)
    tx(s, 6.7, 3.75, 6.0, 1.6, [
        [run("含义:  ", 14.5, True, DARK),
         run("「渲染 → 特征 → 生成模型后验」能端到端反推场景因子。", 14.5, False, DARK)],
        [run("生成模型后验方向可行 —— 触发 08-11 起 SPN 成为主线候选。", 14.5, False, DARK)],
    ], space_after=5)
    footer(s, 9, "97c33d6: N=2000 码 0.970 / kind 1.0 / gx 0.99 / gy 0.98 / size 1.0")
    return s


def s10_retire1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "第一次退役: 离散场景码 → 连续物理量", kicker="两次退役之一",
          evidence="523b97a (08-13) · docs/architecture.md §0")
    shape(s, 0.7, 1.9, 5.6, 3.4, LIGHT, radius=0.06, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
    tx(s, 0.95, 2.1, 5.1, 0.4, [[run("连续化验证 (N=4000, K=64)", 14.5, True, DARK)]])
    tx(s, 0.95, 2.6, 5.2, 2.5, [
        [run("kind 0.897", 22, True, TEAL), run("   白化 1-NN 上限 0.94", 12, False, GRAY)],
        [run("u,v RMSE 6.6px · R²≈0.90", 15, False, DARK)],
        [run("s R² 0.17 · z R² 0.44", 15, False, DARK)],
        [run("s/z 弱是物理: 单目单帧仅乘积可观测(熟悉尺寸歧义)", 11.5, False, GRAY)],
    ], space_after=6)
    shape(s, 6.7, 1.9, 6.0, 3.4, RED_SOFT, line=RED, line_w=1.0,
          radius=0.06, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
    tx(s, 6.95, 2.1, 5.5, 0.4, [[run("退役判据", 14.5, True, RED)]])
    tx(s, 6.95, 2.6, 5.5, 2.5, [
        [run("离散场景码体系(逐码贝叶斯 / 池化 SPN / 码网格)整体退役。", 14, False, DARK)],
        [run("连续物理量的离散化, 只是后验求积 —— 离散化本身不带来任何信息增益。", 14, True, DARK)],
    ], space_after=6)
    tx(s, 0.7, 5.6, 12.0, 0.9, [
        [run("升级路径:  ", 14, True, DARK),
         run("核回归边界饱和不完美(s/z R² 可为负), 后续由实例级 + tied 方差接管。", 14, False, DARK)],
    ], space_after=3)
    footer(s, 10, "523b97a: N=4000 K=64 · kind 0.897 / u,v RMSE 6.6px / s 0.17 / z 0.44")
    return s


def s11_retire2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "第二次退役: EM 质心 → K=N 实例级", kicker="两次退役之二",
          evidence="54cd2d1 · docs/architecture.md §2.2")
    # 对比数字
    shape(s, 0.7, 1.9, 5.6, 2.2, RED_SOFT, line=RED, line_w=1.0, radius=0.06, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
    tx(s, 0.95, 2.05, 5.0, 0.34, [[run("EM / 质心压缩", 14, True, RED)]])
    tx(s, 0.95, 2.5, 5.0, 0.8, [[run("u R² 0.50", 30, True, RED)]])
    tx(s, 0.95, 3.4, 5.0, 0.5, [[run("小数据 + 弯曲流形: 质心把点平均到流形外", 12, False, DARK)]])
    arrow_r(s, 6.45, 2.7, 0.7, 0.4, TEAL)
    shape(s, 7.2, 1.9, 5.5, 2.2, TEAL_SOFT, line=TEAL, line_w=1.0, radius=0.06, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
    tx(s, 7.45, 2.05, 5.0, 0.34, [[run("实例级 (K=N, 无 EM)", 14, True, TEAL)]])
    tx(s, 7.45, 2.5, 5.0, 0.8, [[run("u R² 0.90", 30, True, TEAL)]])
    tx(s, 7.45, 3.4, 5.0, 0.5, [[run("全量 1296: u/v RMSE 5.4/5.2px · R²≈0.93", 12, False, DARK)]])
    # EM 四条退化通道
    tx(s, 0.7, 4.35, 6.0, 0.32, [[run("EM 四条退化通道 (§2.2)", 12.5, True, GRAY)]])
    chans = [
        ("razor 门控", "死分量均值爆炸 90×"),
        ("大方差", "只活 19/216 分量"),
        ("权重死亡", "log_w ≈ −20"),
        ("nk≈3", "方差项 ±115 nats 淹没距离"),
    ]
    x = 0.7
    for i, (t, d) in enumerate(chans):
        gx = 0.7 + i * 3.1
        shape(s, gx, 4.72, 2.9, 1.35, LIGHT, radius=0.1, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
        tx(s, gx + 0.15, 4.85, 2.6, 0.32, [[run(t, 13, True, DARK)]])
        tx(s, gx + 0.15, 5.22, 2.6, 0.75, [[run(d, 11.5, False, GRAY)]])
    tx(s, 0.7, 6.3, 12.0, 0.6, [
        [run("教训:  ", 14.5, True, GOLD),
         run("EM/质心压缩是大数据优化; 数据均匀采样 ⟹ 均匀权重才是正确先验(学权重反而死亡螺旋)。", 14.5, False, DARK)],
    ], space_after=3)
    footer(s, 11, "54cd2d1 · docs/architecture.md §2.2")
    return s


def s12_how(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "当前 MixtureSPN 如何工作", kicker="当前架构",
          evidence="docs/architecture.md §2.1 · §1")
    # 机制修复链 (kind 判别)
    tx(s, 0.7, 1.7, 6.0, 0.32, [[run("kind 判别修复链(逐版实测)", 12.5, True, GRAY)]])
    chain = [("0.47", "无约束 EM 按位置聚类", RED_SOFT),
             ("0.47", "逐 kind 分层拟合", RED_SOFT),
             ("0.47", "色度关 gc + 原始拮抗", RED_SOFT),
             ("0.68", "PCA 白化", BLUE_SOFT),
             ("✓", "方差收缩 → 实例级", TEAL_SOFT)]
    x = 0.7
    for i, (v, t, c) in enumerate(chain):
        gx = 0.7 + i * 2.46
        shape(s, gx, 2.05, 2.26, 1.3, c, radius=0.12, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
        tx(s, gx, 2.15, 2.26, 0.4, [[run(v, 18, True, DARK)]], align=PP_ALIGN.CENTER)
        tx(s, gx + 0.08, 2.62, 2.1, 0.66, [[run(t, 10.5, False, DARK)]], align=PP_ALIGN.CENTER, space_after=1)
        if i < 4:
            arrow_r(s, gx + 2.27, 2.5, 0.17, 0.26, GRAY)
    # 组件结构
    tx(s, 0.7, 3.7, 6.0, 0.32, [[run("实例级组装(无 EM)", 12.5, True, GRAY)]])
    parts = [("逐 kind 分层", "P(kind)·P(f,t|kind)"),
             ("K=N 实例块", "每样本一分量"),
             ("类内 tied 方差", "抑制 nk≈3 噪声"),
             ("均匀权重", "数据均匀采样的先验")]
    x = 0.7
    for i, (t, d) in enumerate(parts):
        gx = 0.7 + i * 3.1
        shape(s, gx, 4.05, 2.9, 1.15, BLUE_SOFT, line=TEAL, line_w=0.75, radius=0.1, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
        tx(s, gx + 0.15, 4.16, 2.6, 0.32, [[run(t, 12.5, True, DARK)]])
        tx(s, gx + 0.15, 4.52, 2.6, 0.6, [[run(d, 10.5, False, GRAY)]])
    # 预测语义
    shape(s, 0.7, 5.5, 12.0, 1.2, LIGHT, radius=0.05, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
    tx(s, 0.95, 5.62, 11.5, 1.0, [
        [run("预测语义:  ", 14, True, DARK),
         run("连续条件期望 E[u,v,s−ŝ,z−ẑ|x] ≡ 分层核回归;  离散因子 P(kind,hue,lcol,ldir|x) = 场景因子后验分类。", 14, False, DARK)],
        [run("白化使对角高斯 ≡ 原空间全协方差 —— 解决「相邻像素强相关 → 对角高斯重复计票」。", 13, False, GRAY)],
    ], space_after=4)
    footer(s, 12, "docs/architecture.md §2.1")
    return s


def s13_bench(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "当前基准: 插值 / 外推 (N=1296, kind_topk=3)", kicker="结果",
          evidence="docs/architecture.md §1")
    # 插值表
    shape(s, 0.7, 1.85, 6.0, 3.3, LIGHT, radius=0.05, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
    tx(s, 0.95, 1.98, 5.5, 0.34, [[run("插值", 16, True, DARK)]])
    rows_i = [("u / v  R²", "0.930 / 0.945"), ("s  R²", "0.508"), ("z  R²", "0.831"),
              ("kind / hue", "0.753 / 1.000"), ("lcol / ldir", "0.994 / 0.895")]
    _metric_rows(s, rows_i, 1.0, 2.42)
    # 外推表
    shape(s, 7.0, 1.85, 5.7, 3.3, LIGHT, radius=0.05, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
    tx(s, 7.25, 1.98, 5.2, 0.34, [[run("外推", 16, True, DARK)]])
    rows_e = [("u / v  R²", "0.949 / 0.953"), ("s / z  R²", "0.922 / 0.956"), ("kind", "0.617"),
              ("hue", "0.981"), ("lcol / ldir", "0.880 / 0.772")]
    _metric_rows(s, rows_e, 7.3, 2.42)
    # 精炼贡献
    shape(s, 0.7, 5.35, 12.0, 1.35, TEAL_SOFT, radius=0.05, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
    tx(s, 0.95, 5.5, 11.5, 0.4, [[run("渲染残差精炼的贡献 (lcol / ldir)", 13.5, True, TEAL)]])
    tx(s, 0.95, 5.95, 11.5, 0.6, [
        [run("0.457 / 0.367", 16, True, RED), run("   →   ", 16, True, GRAY),
         run("0.994 / 0.895", 16, True, TEAL),
         run("   候选重渲染把反照率×光照歧义交回正向模型", 12.5, False, GRAY)],
    ], space_after=3)
    footer(s, 13, "docs/architecture.md §1 · 精炼前 lcol/ldir 0.457/0.367 → 0.994/0.895")
    return s


def _metric_rows(s, rows, x, y):
    for i, (k, v) in enumerate(rows):
        gy = y + i * 0.52
        tx(s, x, gy, 3.2, 0.4, [[run(k, 12, False, GRAY)]])
        tx(s, x + 2.0, gy, 2.8, 0.4, [[run(v, 13.5, True, DARK)]])


def s14_echo(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "方法论回声: 纹理 / roughness 研究", kicker="近期研究",
          evidence="docs/architecture.md §2.3")
    # 左: 谱形轴
    shape(s, 0.7, 1.85, 6.0, 2.6, LIGHT, radius=0.05, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
    tx(s, 0.95, 1.98, 5.5, 0.34, [[run("11 张图不是正交轴, 而是同一 p_s 的聚合", 13.5, True, DARK)]])
    tx(s, 0.95, 2.4, 5.6, 1.9, [
        [run("p_s = e_s / Σe  ", 13, True, TEAL), run("逐像素尺度能量分布", 12, False, GRAY)],
        [run("真正分界: ", 13, True, DARK), run("能量(log_mag) vs 形状(slope/moments)", 13, True, DARK)],
        [run("E7: log_mag 对比度 1.00 · slope/moments 对比度 0.00", 12, False, DARK)],
        [run("slope 对 β 单调(粉 0.73 < 白 1.25 < 蓝 1.28)", 12, False, DARK)],
    ], space_after=5)
    # 右: 纹理/roughness 结论
    shape(s, 7.0, 1.85, 5.7, 2.6, LIGHT, radius=0.05, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
    tx(s, 7.25, 1.98, 5.2, 0.34, [[run("实测: 纹理类型 vs roughness", 13.5, True, DARK)]])
    tx(s, 7.25, 2.4, 5.3, 1.9, [
        [run("纹理类型 ", 13, True, DARK), run("0.499 / 0.452 > chance 0.333", 13, True, TEAL)],
        [run("roughness ", 13, True, DARK), run("二阶弱信号, 被纹理/外观/光照淹没", 13, False, RED)],
        [run("决策: 固定 roughness=0.55, 只保留纹理类型因子", 12.5, False, DARK)],
    ], space_after=5)
    # 底部结论
    shape(s, 0.7, 4.75, 12.0, 1.8, GOLD_SOFT, line=GOLD, line_w=1.0, radius=0.05, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
    tx(s, 0.95, 4.9, 11.5, 1.5, [
        [run("方法论回声:  ", 14.5, True, GOLD),
         run("「实测驱动、诚实否决」再次验证 —— 固定几何探针的乐观(tex 1.00 / roughness 0.997)被主线全变实测修正(0.499 / 负 R²)。", 14.5, False, DARK)],
        [run("特征只暴露数据里的方差, 顺序不可倒: 先给 Codebook 接线自由度, 再立监督目标。", 13.5, False, DARK)],
    ], space_after=5)
    footer(s, 14, "docs/architecture.md §2.3")
    return s


def s15_summary(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "总结: 研究原则与下一步", kicker="结论",
          evidence="docs/architecture.md §7 · §0")
    # 三原则
    prins = [
        ("实测驱动", "每条机制以测量为准, 否决靠数据不靠直觉", TEAL),
        ("诚实否决", "失败证据留档(DGGG/DG-MRF/EM 退化通道)", RED),
        ("分层先验", "硬约束进支持集 · 统计进概率 · 歧义进候选后验", GOLD),
    ]
    x = 0.7
    for i, (t, d, c) in enumerate(prins):
        gx = 0.7 + i * 4.1
        shape(s, gx, 1.85, 3.85, 1.5, LIGHT, radius=0.08, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
        shape(s, gx, 1.85, 0.12, 1.5, c)
        tx(s, gx + 0.25, 1.98, 3.4, 0.34, [[run(t, 15, True, c)]])
        tx(s, gx + 0.25, 2.36, 3.45, 0.9, [[run(d, 12.5, False, DARK)]])
    # 核心结论
    tx(s, 0.7, 3.7, 12.0, 0.34, [[run("核心结论", 12.5, True, GRAY)]])
    shape(s, 0.7, 4.05, 12.0, 1.35, DARK, radius=0.05, kind=MSO_SHAPE.ROUNDED_RECTANGLE)
    tx(s, 0.95, 4.2, 11.5, 1.1, [
        [run("显式建模的正确粒度 = 生成过程, 而非推理模块。", 16, True, WHITE)],
        [run("小数据 + 弯曲流形 → 实例级(K=N, 无 EM); 数据均匀采样 → 均匀权重是正确先验。", 13, False, RGBColor(0xC8, 0xD2, 0xE0))],
    ], space_after=4)
    # 下一步
    tx(s, 0.7, 5.65, 12.0, 0.34, [[run("下一步", 12.5, True, GRAY)]])
    nxt = ["双层渲染残差", "池外光照探针", "逐 kind PPCA 似然比", "大数据逃生通道"]
    x = 0.7
    for i, t in enumerate(nxt):
        gx = 0.7 + i * 3.1
        chip(s, gx, 6.05, 2.9, 0.52, t, TEAL, size=12)
    footer(s, 15, "docs/architecture.md §7 · §0")
    return s


# ───────────────────────────────────────────────────────────────────────────
# 构建 + 校验
# ───────────────────────────────────────────────────────────────────────────

SLIDES = [s1_title, s2_overview, s3_problem, s4_pipeline, s5_old_pipeline,
          s6_kept, s7_fail1, s8_fail2, s9_pivot, s10_retire1, s11_retire2,
          s12_how, s13_bench, s14_echo, s15_summary]


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for fn in SLIDES:
        fn(prs)
    return prs


def _est_text_overflow(shapes_):
    """返回 (警告列表)。用字符宽度启发式估算每文本框是否溢出。"""
    warns = []
    for sp in shapes_:
        if not sp.has_text_frame:
            continue
        tf = sp.text_frame
        w_in = sp.width / 914400
        h_in = sp.height / 914400
        total = 0.0
        for p in tf.paragraphs:
            line_w = 0.0
            line_h = 0.0
            for r in p.runs:
                size = (r.font.size.pt if r.font.size else 14) / 72.0
                for ch in r.text:
                    cw = _char_w_in(ch, size)
                    if line_w + cw > w_in and line_w > 0:
                        total += line_h * 1.16
                        line_w = cw
                        line_h = size
                    else:
                        line_w += cw
                        line_h = max(line_h, size)
            total += line_h * 1.16
        if total > h_in * 1.03:
            txt = "".join(r.text for p in tf.paragraphs for r in p.runs)[:24]
            warns.append(f"    overflow? box {w_in:.1f}x{h_in:.1f}in est {total:.2f}in :: '{txt}'")
    return warns


def validate(path: Path) -> dict:
    prs = Presentation(str(path))
    n = len(prs.slides)
    errs = []
    warns = []
    if n != TOTAL:
        errs.append(f"slide count {n} != {TOTAL}")
    w_in = prs.slide_width / 914400
    h_in = prs.slide_height / 914400
    if abs(w_in - SLIDE_W_IN) > 0.01 or abs(h_in - SLIDE_H_IN) > 0.01:
        errs.append(f"size {w_in:.3f}x{h_in:.3f} != 13.333x7.5")
    titles = []
    for i, sl in enumerate(prs.slides, 1):
        ti = None
        for sp in sl.shapes:
            l, t = sp.left / 914400, sp.top / 914400
            w, h = sp.width / 914400, sp.height / 914400
            if l < -0.01 or t < -0.01 or l + w > SLIDE_W_IN + 0.01 or t + h > SLIDE_H_IN + 0.01:
                errs.append(f"slide {i}: shape out of bounds ({l:.2f},{t:.2f},{w:.2f},{h:.2f})")
            if sp.has_text_frame and ti is None:
                txt = "".join(r.text for p in sp.text_frame.paragraphs for r in p.runs).strip()
                if len(txt) > 4 and not txt.startswith("CONGER"):
                    ti = txt[:40]
        for wmsg in _est_text_overflow(sl.shapes):
            warns.append(f"slide {i}: {wmsg}")
        titles.append(ti or "(blank)")
    return {"count": n, "size_in": f"{w_in:.3f}x{h_in:.3f}", "errors": errs,
            "warnings": warns, "titles": titles}


def main() -> None:
    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(str(OUT_PPTX))
    print(f"[ok] wrote {OUT_PPTX}")
    res = validate(OUT_PPTX)
    print(f"[validate] slides={res['count']}  size={res['size_in']}")
    for i, t in enumerate(res["titles"], 1):
        print(f"  {i:2d}. {t}")
    if res["errors"]:
        print("\n[validate] ERRORS:")
        for e in res["errors"]:
            print("  !!", e)
    else:
        print("\n[validate] bounds/geometry: OK")
    if res["warnings"]:
        print("[validate] overflow warnings (heuristic):")
        for w in res["warnings"]:
            print(w)
    else:
        print("[validate] overflow heuristic: no warnings")
    print(f"[validate] {'FAIL' if res['errors'] else 'PASS'}")


if __name__ == "__main__":
    main()
